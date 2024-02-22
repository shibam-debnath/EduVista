import os
import io
import base64
from gtts import gTTS
import comtypes.client
from io import BytesIO
from pptx import Presentation
from dotenv import load_dotenv
import google.generativeai as genai

from flask_bcrypt import Bcrypt
from flask_cors import CORS, cross_origin
from flask import Flask, request, jsonify, send_file, request, make_response
from flask_jwt_extended import create_access_token,get_jwt,get_jwt_identity, unset_jwt_cookies, jwt_required, JWTManager

from models import db, User

# Create a Flask app
app = Flask(__name__)
CORS(app)
load_dotenv()

app.config['SECRET_KEY'] = 'h4saki-secret-key'
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///flaskdb.db'

SQLALCHEMY_TRACK_MODIFICATIONS = False
SQLALCHEMY_ECHO = True


jwt = JWTManager(app)
bcrypt = Bcrypt(app)   
db.init_app(app) 

with app.app_context():
    db.create_all()



GOOGLE_API_KEY = "AIzaSyAa7UYKqvP_17wdthjVwgz4RPbJIRnro8M"
genai.configure(api_key = GOOGLE_API_KEY)
model = genai.GenerativeModel('gemini-pro')



@app.route('/logintoken', methods=["POST"])
def create_token():
    email = request.json.get("email", None)
    password = request.json.get("password", None)

    user = User.query.filter_by(email=email).first()

    if user is None:
        return jsonify({"error": "Wrong email or passwords"}), 401
    
    # compare passwords
    if not bcrypt.check_password_hash(user.password, password):
        return jsonify({"error": "Unauthorized"}), 401

    access_token = create_access_token(identity=email)
    return jsonify({
        "email": email,
        "access_token": access_token
    })

@app.route("/logout", methods=["POST"])
def logout():
    response = jsonify({"msg": "logout successful"})
    unset_jwt_cookies(response)
    return response



@app.route('/signup', methods=["POST", "OPTIONS"])
def signup():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()

    if request.method == "POST":
        name = request.json.get("name", None)
        email = request.json.get("email", None)
        password = request.json.get("password", None)

        user_exists = User.query.filter_by(email=email).first() is not None
    
        if user_exists:
            return jsonify({"error": "Email already exists"}), 409
        
        hashed_password = bcrypt.generate_password_hash(password)
        new_user = User(name=name, email=email, password=hashed_password, about="sample about me")
        db.session.add(new_user)
        db.session.commit()
    
        return jsonify({
            "id": new_user.id,
            "email": new_user.email,
        })


@app.route("/")
@cross_origin(supports_credentials=True)
def home():
    return jsonify("Hello Shibam! Flask Server is running properly")


# Define the directory where you want to store uploaded PPTX files
UPLOAD_FOLDER = 'static'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# function that extracts file information from the ppt using python-pptx
def read_ppt(file_path):

    presentation = Presentation(file_path)

    slide_headings = {}
    slide_contents = {}
    all_slide_contents = ""

    for slide_number, slide in enumerate(presentation.slides):
        heading = f"Slide {slide_number + 1}"

        text = ""
        heading_text = ""

        for shape_number, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                if heading_text == "" and shape.text.strip() != "":
                    heading_text = shape.text.strip()
                else:
                    text += shape.text + " "
                    all_slide_contents += shape.text + " "

        slide_headings[heading] = heading_text
        slide_contents[heading] = text.strip()


    return slide_headings, slide_contents, all_slide_contents.strip()



def convert_pptx_to_pdf(pptx_path, output_dir):

    try:
        # Ensure output directory exists
        os.makedirs(output_dir, exist_ok=True)

        # Get filename without extension
        filename = os.path.splitext(os.path.basename(pptx_path))[0]
        output_file_path = os.path.join(output_dir, f"{filename}.pdf")
        return output_file_path



        print(output_file_path)

        # Create PowerPoint application object
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")


        # Minimize PowerPoint window (optional)
        powerpoint.Visible = 0  # Set to 1 to make it visible

        # Open the PPTX file
        slides = powerpoint.Presentations.Open(pptx_path)

        # Save as PDF
        slides.SaveAs(output_file_path, 32)  # formatType = 32 for PDF

        # Close PowerPoint objects
        slides.Close()
        powerpoint.Quit()

        return output_file_path

    except FileNotFoundError as e:
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}") from e
    except OSError as e:
        raise OSError(f"Error saving PDF file: {output_file_path}") from e



@app.route('/upload', methods=['POST', "OPTIONS"])
def upload_file():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    
    if request.method == "POST":
        if 'file' not in request.files:
            return jsonify({'error': 'No file part'})

        file = request.files['file']
        avatar = request.form.get('avatar')

        if file and file.filename.endswith('.pptx'):
            
            os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

            base_file_path = os.path.join(app.config['UPLOAD_FOLDER'])
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], "slides.pptx")
            file.save(file_path)


            # Initialize COM library
            comtypes.CoInitialize()

            # Get console arguments
            input_file_path = "static/slides.pptx"
            output_file_path = "static/output.pdf"

            # Convert file paths to Windows format
            input_file_path = os.path.abspath(input_file_path)
            output_file_path = os.path.abspath(output_file_path)

            # Create PowerPoint application object
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")

            # Set visibility to minimize
            powerpoint.Visible = 1

            # Open the PowerPoint slides
            slides = powerpoint.Presentations.Open(input_file_path)

            # Save as PDF (formatType = 32)
            slides.SaveAs(output_file_path, 32)

            # Close the slide deck
            slides.Close()

            pdf_path ="static/output.pdf"
            return send_file(pdf_path, as_attachment=True)

        else: 
            return jsonify({'error': 'Invalid file'})
    else:
        raise RuntimeError("Weird - don't know how to handle method {}".format(request.method))









# API 0 : model 0 ( base model ) : reads the ppt 
@app.route("/api/generate_content", methods=["POST", "OPTIONS"])
def generate_content():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    if request.method == "POST":
        if 'file' not in request.files:
            return jsonify({'error': 'No file part'})

        file = request.files['file']
        avatar = request.form.get('avatar')

        if file and file.filename.endswith('.pptx'):
            
            os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

            file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(file_path)

            headings, contents, all_slide_contents = read_ppt(file_path)
            all_slide_contents = all_slide_contents.replace('*', '').strip()
            
            print(f"All slide contents : {all_slide_contents}")

            return jsonify({
                            
                'message': 'Content generated Successfully', 
                'headings': headings, 
                'contents': contents,
                'all_slide_contents': all_slide_contents, 
                'avatar': avatar

            })

        else: 
            return jsonify({'error': 'Invalid file'})
    else:
        raise RuntimeError("Weird - don't know how to handle method {}".format(request.method))




@app.route('/api/generate_content_gemini', methods=["POST", "OPTIONS"])
def generate_content_gemini():

    if request.method == "OPTIONS":
        return _build_cors_preflight_response()

    if request.method == "POST":

        if 'file' not in request.files:
            return jsonify({'error': 'No file part'})

        file = request.files['file']
        avatar = request.form.get('avatar')

        if file and file.filename.endswith('.pptx'):
            
            os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

            file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(file_path)

            headings, contents, all_slide_contents = read_ppt(file_path)

            prompt = f"Hey gemini I am giving you the contens of a ppt and you need to Generate a content/script from the prospective of a teacher that descripes the ppt contents within 200 words. Here is the ppt content : {all_slide_contents}"
            res = model.generate_content(prompt)
            res = res.text
            formatted_response = res.replace('*', '').strip()

            return jsonify({
                            
                'message': 'Content generated Successfully', 
                'headings': headings, 
                'contents': contents,
                'all_slide_contents': formatted_response, 
                'avatar': avatar
                
            })


    else:
        raise RuntimeError("Weird - don't know how to handle method {}".format(request.method))


# API 4 : generate text to audio
@app.route("/api/generate_audio", methods=["POST", "OPTIONS"])
def generate_audio():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()

    if request.method == "POST":
        data = request.get_json()
        text = data['text']
        voice = data.get('voice', 'en')  # default to English

        tts = gTTS(text=text, lang=voice, slow=False)
        audio_file = io.BytesIO()
        tts.write_to_fp(audio_file)

        # Reset file pointer to the beginning
        audio_file.seek(0)

        return send_file(audio_file, mimetype='audio/mpeg', as_attachment=True, download_name='generated_audio.mp3')
    else:
        raise RuntimeError("Weird - don't know how to handle method {}".format(request.method))



# handling CORS preflight 
def _build_cors_preflight_response():
    response = make_response()
    response.headers.add("Access-Control-Allow-Origin", "*")
    response.headers.add('Access-Control-Allow-Headers', "*")
    response.headers.add('Access-Control-Allow-Methods', "*")
    return response


# main 
if __name__ == "__main__":
    app.run(debug=True)

