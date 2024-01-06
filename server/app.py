import os
import io
import time
import torch
import openai
import locale
import requests
import transformers
from gtts import gTTS
from pptx import Presentation
from openai import Completion
from pydub import AudioSegment
from IPython.display import Audio
from langchain.chains import LLMChain
from transformers import AutoTokenizer
from flask_cors import CORS, cross_origin
from langchain.prompts import PromptTemplate
from langchain.llms import HuggingFacePipeline
from flask import Flask, request, jsonify, send_file, request, make_response

import torch
from transformers import GPT2Tokenizer, GPT2LMHeadModel
from flask import Flask, request, jsonify


# Create a Flask app
app = Flask(__name__)
CORS(app)


# Load GPT-2 tokenizer and model
tokenizer = GPT2Tokenizer.from_pretrained('gpt2')
model = GPT2LMHeadModel.from_pretrained('gpt2')


# API key for chatGPT 3.5
openai.api_key = "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"

# Define the directory where you want to store uploaded PPTX files
UPLOAD_FOLDER = 'static/pptx_files'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# function that extracts file information from the ppt using python-pptx
def read_ppt(file_path):

    presentation = Presentation(file_path)

    slide_headings = {}
    slide_contents = {}
    all_slide_contents = ""

    for slide_number, slide in enumerate(presentation.slides):
        heading = f"Slide {slide_number + 1}"

        text = ""
        heading_text = ""

        for shape_number, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                if heading_text == "" and shape.text.strip() != "":
                    heading_text = shape.text.strip()
                else:
                    text += shape.text + " "
                    all_slide_contents += shape.text + " "

        slide_headings[heading] = heading_text
        slide_contents[heading] = text.strip()


    return slide_headings, slide_contents, all_slide_contents.strip()



# API 1 : home API to check if everything is working fine or not 
@app.route("/")
@cross_origin(supports_credentials=True)
def home():
    return jsonify("Hello Shibam! Flask Server is running properly")



# API 2 : generate script for using the model

# model = "meta-llama/Llama-2-7b-chat-hf"
# tokenizer = AutoTokenizer.from_pretrained(model)

# pipeline = transformers.pipeline(
#     "text-generation",
#     model=model,
#     tokenizer=tokenizer,
#     torch_dtype=torch.bfloat16,
#     trust_remote_code=True,
#     device_map="auto",
#     max_length=500,
#     do_sample=True,
#     top_k=10,
#     num_return_sequences=1,
#     eos_token_id=tokenizer.eos_token_id
# )

# llm = HuggingFacePipeline(pipeline=pipeline, model_kwargs={'temperature': 0})

# prompt_template = PromptTemplate(
#     input_variables=["text", "heading", "title"],
#     template="Explain the each points in {text} for {heading} slide in the context of {title}"
# )

# chain = LLMChain(llm=llm, prompt=prompt_template, verbose=True)


# @app.route("/api/generate_content", methods=["POST", "OPTIONS"])
# def generate_content():
#     # Get the request data
#     if request.method == "OPTIONS":
#         return _build_cors_preflight_response()
#     else:
#         return _build_cors_preflight_response()

    # if request.method == "POST":
        # Check if the POST request has the file part
    #     if 'file' not in request.files:
    #         return jsonify({'error': 'No file part'})

    #     file = request.files['file']
    #     avatar = request.form.get('avatar')

    #     # Check if the file is present and has an allowed file extension
    #     if file and file.filename.endswith('.pptx'):
    #         # Ensure the upload folder exists
    #         os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

    #         # Save the file to the configured upload folder
    #         file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
    #         file.save(file_path)

    #         # Read the ppt file to get headings, contents, and all_slide_contents
    #         headings, contents, all_slide_contents = read_ppt(file_path)
    #         title="AI powered Voice guided PPT explanation"

    #         # Modify this part based on your needs
    #         generated_text = {}
    #         content_text = ""
    #         for (slide_head, heading), (slide_content, content) in zip(headings.items(), contents.items()):
    #             response = chain.run({"text": content, "heading": heading, "title": title})
    #             generated_text[slide_head] = response
    #             content_text += response + " "

    #         # Optionally, you can send the extracted content to the frontend
    #         return jsonify({"generated_text": generated_text, "content_text": content_text.strip()})

    #     return jsonify({'error': 'Invalid file'})
    # else:
    #     raise RuntimeError("Weird - don't know how to handle method {}".format(request.method))





def generate_content_with_chatgpt(prompt):
    try:
        # Get PowerPoint content from the request
        ppt_content = prompt


        # Use ChatGPT-3.5 API to generate a script
        response = openai.Completion.create(
            engine="text-davinci-003",  # You can experiment with different engines
            # prompt=f"Generate a teaching script for the following PowerPoint content:\n\n{ppt_content}",
            prompt=f"who are you ?",
            max_tokens=500  # Adjust as needed
        )

        # Extract the generated script from the API response
        generated_script = response['choices'][0]['text']
        print("generatd script : ")
        print(generated_script)

        return generated_script

    except Exception as e:
        print(f"Error: {str(e)}")
        return "error"



@app.route("/api/generate_content", methods=["POST", "OPTIONS"])
def generate_content():
    # Get the request data
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()
    if request.method == "POST":
        # Check if the POST request has the file part
        if 'file' not in request.files:
            return jsonify({'error': 'No file part'})

        file = request.files['file']
        avatar = request.form.get('avatar')

        # Check if the file is present and has an allowed file extension
        if file and file.filename.endswith('.pptx'):
            # Ensure the upload folder exists
            os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

            # Save the file to the configured upload folder
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(file_path)

            # Read the ppt file to get headings, contents, and all_slide_contents
            headings, contents, all_slide_contents = read_ppt(file_path)

            
            # chatGPT

            # return jsonify({"upar hi dediya": all_slide_contents.strip()})
            # content_text = generate_content_with_chatgpt(all_slide_contents)
            # # Now you can directly return the generated script
            # return jsonify({"content_text": content_text.strip()})

            # Tokenize input text
            input_ids = tokenizer.encode(all_slide_contents, return_tensors="pt")
            # Generate script using GPT-2
            outputs = model.generate(input_ids, max_length=500, num_beams=5, no_repeat_ngram_size=2)
            # Decode the generated script
            generated_script = tokenizer.decode(outputs[0], skip_special_tokens=True)
            return jsonify({'script': generated_script})

            # Modify this part based on your needs
            # content_text = ""
            # generated_text = {}
            # for (slide_head, heading), (slide_content, content) in zip(headings.items(), contents.items()):
            #     # Use ChatGPT API to generate content based on the slide content
            #     chatgpt_response = generate_content_with_chatgpt(content)

            #     # Modify this line based on your concatenation requirement
            #     content_text += chatgpt_response.strip() + " "
            #     generated_text[slide_head] = chatgpt_response

            # return jsonify({"generated_text": generated_text, "all_slide_contents": all_slide_contents,
            #                 "content_text": content_text.strip()})

        else: 
            return jsonify({'error': 'Invalid file'})
    else:
        raise RuntimeError("Weird - don't know how to handle method {}".format(request.method))


# API 3 : model 0 ( base model ) on the `/upload` route
@app.route('/upload', methods=['POST', "OPTIONS"])
def upload_file():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()

    if request.method == "POST":
        # Check if the POST request has the file 
        if 'file' not in request.files:
            return jsonify({'error': 'No file'})

        file = request.files['file']
        avatar = request.form.get('avatar')

        # Check if the file is present and has an allowed file extension
        if file and file.filename.endswith('.pptx'):
            # Ensure the upload folder exists
            os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

            # Save the file to the configured upload folder
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(file_path)

            # Read PPT and extract text
            headings, contents, all_slide_contents  = read_ppt(file_path)

            # Optionally, you can send the extracted content to the frontend
            return jsonify({'message': 'File uploaded successfully', 'headings': headings, 'contents': contents,'all_slide_contents': all_slide_contents, 'avatar': avatar})

        return jsonify({'error': 'Invalid file'})
    else:
        raise RuntimeError("Weird - don't know how to handle method {}".format(request.method))




# API 4 : generate text to audio
@app.route("/api/generate_audio", methods=["POST", "OPTIONS"])
def generate_audio():
    if request.method == "OPTIONS":
        return _build_cors_preflight_response()

    if request.method == "POST":
        data = request.get_json()
        text = data['text']
        voice = data.get('voice', 'en')  # default to English

        # Generate audio using gTTS
        tts = gTTS(text=text, lang=voice, slow=False)
        audio_file = io.BytesIO()
        tts.write_to_fp(audio_file)

        # Reset file pointer to the beginning
        audio_file.seek(0)

        # Send the audio file directly
        return send_file(audio_file, mimetype='audio/mpeg', as_attachment=True, download_name='generated_audio.mp3')
    else:
        raise RuntimeError("Weird - don't know how to handle method {}".format(request.method))







# handling CORS preflight 
def _build_cors_preflight_response():
    response = make_response()
    response.headers.add("Access-Control-Allow-Origin", "*")
    response.headers.add('Access-Control-Allow-Headers', "*")
    response.headers.add('Access-Control-Allow-Methods', "*")
    return response
def _corsify_actual_response(response):
    response.headers.add("Access-Control-Allow-Origin", "*")
    return response

# main 
if __name__ == "__main__":
    app.run(debug=True)
